#!/usr/bin/env python3
# /// script
# dependencies = ["python-pptx>=1.0", "Pillow>=10"]
# ///
"""Generate Ailurus Overflow demo deck aligned with hackathon video structure."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
LOGO = ROOT / "ailurus-web" / "public" / "logo.png"
OUT = ROOT / "docs" / "ailurus-overflow-demo-deck.pptx"

INK = RGBColor(0x1a, 0x1a, 0x1a)
MUTED = RGBColor(0x6b, 0x6b, 0x6b)
PANDA = RGBColor(0xe8, 0x5d, 0x3c)
CREAM = RGBColor(0xf7, 0xf3, 0xed)
WHITE = RGBColor(0xff, 0xff, 0xff)


def set_slide_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title: str, subtitle: str = "", timing: str = ""):
    if timing:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(3), Inches(0.4))
        p = box.text_frame.paragraphs[0]
        p.text = timing
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = PANDA

    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.75), Inches(12), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = INK

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = MUTED
        p2.space_before = Pt(12)


def add_bullets(slide, items: list[str], top=2.0):
    box = slide.shapes.add_textbox(Inches(0.75), Inches(top), Inches(11.5), Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(22)
        p.font.color.rgb = INK
        p.space_after = Pt(14)


def add_footer(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(6.85), Inches(12), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = MUTED


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 Cover
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, CREAM)
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(1.2), Inches(1.4), width=Inches(4.2))
    box = slide.shapes.add_textbox(Inches(1.2), Inches(3.5), Inches(10), Inches(2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Creator subscriptions on Sui"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = INK
    p2 = tf.add_paragraph()
    p2.text = "Sui Overflow 2026 · Walrus Track · Testnet MVP"
    p2.font.size = Pt(18)
    p2.font.color.rgb = MUTED
    p2.space_before = Pt(10)
    add_footer(slide, "Google login · USDC · Walrus · Seal · Zero gas")

    # 2 Problem overview (30–60s)
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)
    add_title(
        slide,
        "The problem",
        "Who hurts today, and why existing tools fail",
        "30–60 sec · Problem overview",
    )
    add_bullets(
        slide,
        [
            "Creators lose 20%+ to platforms and do not own their files",
            "Fans want subscriptions, not wallets, gas, or seed phrases",
            "Centralized storage means deplatforming and opaque payouts",
            "Cross-border fans need global stablecoin access, not card rails",
        ],
    )
    add_footer(slide, "Market validated by Patreon / OnlyFans — UX and ownership still broken")

    # 3 Solution snapshot
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, CREAM)
    add_title(slide, "Ailurus in one line", "Web2 experience, full Sui stack underneath")
    add_bullets(
        slide,
        [
            "Fans: Google login → USDC subscribe → instant unlock",
            "Creators: upload to Walrus, price in USDC, earn on-chain",
            "Invisible: Enoki zkLogin + sponsored gas + Seal paywalls",
        ],
        top=2.2,
    )

    # 4 Demo intro (~3 min)
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)
    add_title(
        slide,
        "Product demo",
        "Live Testnet walkthrough",
        "~3 min · Product demo",
    )
    add_bullets(
        slide,
        [
            "1. Landing — publish privately, get paid globally",
            "2. Feed — browse on-chain creator posts",
            "3. Explore — discover profiles and subscription prices",
            "4. Subscribe — USDC payment, Move contract records access",
            "5. Unlock — Seal decrypts media for active subscribers",
            "6. Upload — Walrus SDK + encrypted drops for creators",
        ],
    )

    # 5 Sui stack
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, INK)
    box = slide.shapes.add_textbox(Inches(0.75), Inches(0.75), Inches(11), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = "Full Sui stack — not a bolt-on demo"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    items = [
        "Enoki — Google identity + sponsored gas (Worker backend)",
        "Move — USDC subscriptions, profiles, post metadata on Testnet",
        "Walrus — encrypted albums and clips at scale",
        "Seal — subscription-gated decryption policies",
        "Cetus — USDC→WAL liquidity for storage (user-invisible)",
    ]
    box2 = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11), Inches(4))
    tf = box2.text_frame
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(0xe8, 0xe8, 0xe8)
        p.space_after = Pt(16)

    # 6 Evidence
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)
    add_title(slide, "On-chain evidence (Testnet)", "Deployed and verifiable today")
    add_bullets(
        slide,
        [
            "Package: 0xe5f702…55937",
            "Platform: 0xf6f9d5…3cada",
            "Circle testnet USDC subscriptions",
            "GitHub: github.com/jeffierw/ailurus",
        ],
        top=2.1,
    )

    # 7 Conclusion + vision (30–60s)
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, CREAM)
    add_title(
        slide,
        "Conclusion & vision",
        "Production-ready path beyond the hackathon",
        "30–60 sec · Future vision",
    )
    add_bullets(
        slide,
        [
            "Ailurus = Web2-grade creator subscriptions on decentralized infrastructure",
            "Aha moment: login → pay USDC → unlock encrypted content in seconds",
            "Mainnet target: July 2026 — contracts, Enoki prod, onboarding",
            "Next: multi-OAuth, analytics, recommendation feed, Enoki Connect",
        ],
        top=2.1,
    )

    # 8 Thank you
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, CREAM)
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(4.8), Inches(1.8), width=Inches(3.5))
    box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.3), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Thank you · Q&A"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = INK
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = "ailurus · Sui Overflow 2026 · Walrus Track"
    p2.font.size = Pt(18)
    p2.font.color.rgb = MUTED
    p2.space_before = Pt(12)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Saved {OUT}")


if __name__ == "__main__":
    main()
